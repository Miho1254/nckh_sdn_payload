#!/home/miho/Downloads/nckh/Presentation_Package/.venv/bin/python3
"""
Export HTML Presentation to PowerPoint (PPTX)
Sử dụng Playwright để chụp screenshot từng slide và python-pptx để tạo PPTX.
"""

import os
import sys
import time
from pathlib import Path

# Add venv site-packages to path BEFORE any other imports
sys.path.insert(0, '/home/miho/Downloads/nckh/Presentation_Package/.venv/lib/python3.14/site-packages')

# Now import dependencies
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt

# Configuration
HTML_FILE = Path(__file__).parent / "index.html"
OUTPUT_DIR = Path(__file__).parent / "slides_export"
OUTPUT_PPTX = Path(__file__).parent / "NCKH_TFTDQN_Presentation.pptx"
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720
TOTAL_SLIDES = 14  # Slide 0-13


def capture_slides():
    """Capture screenshots of all slides using Playwright."""
    print("📸 Đang khởi động trình duyệt...")
    
    OUTPUT_DIR.mkdir(exist_ok=True)
    
    with sync_playwright() as p:
        # Launch browser
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
            device_scale_factor=2  # High resolution
        )
        page = context.new_page()
        
        # Load HTML file
        html_url = f"file://{HTML_FILE.absolute()}"
        print(f"📄 Đang mở: {html_url}")
        page.goto(html_url, wait_until="networkidle")
        time.sleep(1)  # Wait for initial render
        
        # Capture each slide
        for slide_idx in range(TOTAL_SLIDES):
            print(f"  🖼️  Chụp slide {slide_idx}/{TOTAL_SLIDES - 1}...")
            
            # Navigate to slide using JavaScript
            page.evaluate(f"window.showSlide({slide_idx})")
            time.sleep(0.5)  # Wait for animation
            
            # Take screenshot
            slide_path = OUTPUT_DIR / f"slide_{slide_idx:02d}.png"
            page.screenshot(path=str(slide_path), full_page=False)
        
        browser.close()
    
    print(f"✅ Đã chụp {TOTAL_SLIDES} slides tại: {OUTPUT_DIR}")
    return OUTPUT_DIR


def create_presentation():
    """Create PowerPoint presentation from screenshots."""
    print("\n📊 Đang tạo file PowerPoint...")
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 ratio
    prs.slide_height = Inches(7.5)
    
    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # Add each slide
    for slide_idx in range(TOTAL_SLIDES):
        slide_path = OUTPUT_DIR / f"slide_{slide_idx:02d}.png"
        
        if not slide_path.exists():
            print(f"  ⚠️  Bỏ qua slide {slide_idx} (không tìm thấy ảnh)")
            continue
        
        print(f"  ➕ Thêm slide {slide_idx}...")
        
        # Add blank slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Add image to fill the entire slide
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        
        slide.shapes.add_picture(str(slide_path), left, top, width, height)
    
    # Save presentation
    prs.save(str(OUTPUT_PPTX))
    print(f"\n✅ Đã tạo file PowerPoint: {OUTPUT_PPTX}")
    return OUTPUT_PPTX


def main():
    print("=" * 60)
    print("🎬 EXPORT HTML PRESENTATION TO POWERPOINT")
    print("=" * 60)
    
    if not HTML_FILE.exists():
        print(f"❌ Không tìm thấy file HTML: {HTML_FILE}")
        sys.exit(1)
    
    # Step 1: Capture slides
    capture_slides()
    
    # Step 2: Create PowerPoint
    pptx_path = create_presentation()
    
    print("\n" + "=" * 60)
    print("🎉 HOÀN TẤT!")
    print(f"📁 File PowerPoint: {pptx_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()